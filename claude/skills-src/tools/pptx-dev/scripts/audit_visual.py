#!/usr/bin/env python3
"""Inspect a .pptx and emit visual-audit findings as JSON."""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.util import Emu

sys.path.insert(0, str(Path(__file__).resolve().parent))

from helpers.render import contrast_ratio  # noqa: E402


SLIDE_WIDTH_EMU = 12192000
SLIDE_HEIGHT_EMU = 6858000

PT_TO_EMU = 12700  # 1pt = 12700 EMU
MIN_FONT_BODY_PT = 14
MIN_FONT_TITLE_PT = 24
MIN_CONTRAST = 4.5
MIN_MARGIN_PT = 24
MAX_ALIGN_DRIFT_PT = 2


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Visual audit of a .pptx")
    p.add_argument("--pptx", type=Path, required=True)
    p.add_argument("--tokens", type=Path, required=True)
    p.add_argument("--output", type=Path, required=True)
    return p.parse_args()


def emu_to_pt(emu: int) -> float:
    return emu / PT_TO_EMU


def _own_fill_hex(fill) -> str | None:
    """Return a FillFormat's own solid color as #RRGGBB, or None (no explicit
    solid fill set / inherited / patterned / gradient) — callers should fall
    back to the slide/page background in that case."""
    try:
        if fill.type == MSO_FILL.SOLID:
            return f"#{str(fill.fore_color.rgb)}"
    except TypeError:
        pass  # non-solid FillFormat (_NoneFill / _NoFill / _GradFill, ...) has no .rgb
    except AttributeError:
        pass  # solid but theme-color based (_SchemeColor) has no .rgb, only .theme_color
    return None


def collect_text_frames(slide):
    """Yield (shape, text_frame, own_bg_hex) for shapes that have text.

    Includes table cells: a GraphicFrame holding a table has has_text_frame=False
    (python-pptx exposes cell text via shape.table.rows[i].cells[j].text_frame,
    a different path), so table content would otherwise be invisible to every
    check that relies on this helper (font size, contrast, empty-slide).

    own_bg_hex is the shape's (or table cell's) own solid fill color, used by
    contrast checks instead of blindly assuming the slide/page background —
    a table header with light text is normally set against a dark cell fill,
    not the page background.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape, shape.text_frame, _own_fill_hex(shape.fill)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield shape, cell.text_frame, _own_fill_hex(cell.fill)


def is_title_like(shape) -> bool:
    """Heuristic: shape sitting in the top ~15% of the slide is treated as title."""
    return shape.top is not None and shape.top < SLIDE_HEIGHT_EMU * 0.15


def audit_fonts(prs, tokens_theme: dict) -> list[dict]:
    body_min = tokens_theme.get("font_size_body_min", MIN_FONT_BODY_PT)
    title_min = tokens_theme.get("font_size_title_min", MIN_FONT_TITLE_PT)
    violations: list[dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape, tf, _own_bg in collect_text_frames(slide):
            min_pt = title_min if is_title_like(shape) else body_min
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size is None:
                        continue
                    pt = run.font.size.pt
                    if pt < min_pt:
                        violations.append({
                            "slide": slide_idx,
                            "shape": shape.name,
                            "text": run.text[:40],
                            "size_pt": pt,
                            "min_pt": min_pt,
                        })
    return violations


def audit_contrast(prs, theme: dict, plan_layouts: dict[int, str]) -> list[dict]:
    default_bg = theme.get("color_bg", "#FFFFFF")
    section_bg = theme.get("color_section_bg", default_bg)
    violations: list[dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_bg = section_bg if plan_layouts.get(slide_idx) == "section" else default_bg
        for shape, tf, own_bg in collect_text_frames(slide):
            bg = own_bg if own_bg is not None else slide_bg
            for para in tf.paragraphs:
                for run in para.runs:
                    color = run.font.color
                    if color is None or color.type is None or color.rgb is None:
                        continue
                    fg = f"#{str(color.rgb)}"
                    ratio = contrast_ratio(fg, bg)
                    if ratio < MIN_CONTRAST:
                        violations.append({
                            "slide": slide_idx,
                            "shape": shape.name,
                            "text": run.text[:40],
                            "fg": fg,
                            "bg": bg,
                            "ratio": round(ratio, 2),
                            "min_ratio": MIN_CONTRAST,
                        })
    return violations


def audit_overflow(prs) -> list[dict]:
    """A shape overflows if its right/bottom edge exceeds slide bounds."""
    violations: list[dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            right = (shape.left or 0) + (shape.width or 0)
            bottom = (shape.top or 0) + (shape.height or 0)
            if right > SLIDE_WIDTH_EMU or bottom > SLIDE_HEIGHT_EMU:
                violations.append({
                    "slide": slide_idx,
                    "shape": shape.name,
                    "overflow_right_pt": round(emu_to_pt(max(0, right - SLIDE_WIDTH_EMU)), 1),
                    "overflow_bottom_pt": round(emu_to_pt(max(0, bottom - SLIDE_HEIGHT_EMU)), 1),
                })
    return violations


def audit_alignment(prs) -> list[dict]:
    """Detect drift > MAX_ALIGN_DRIFT_PT between similarly-positioned shapes."""
    violations: list[dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        lefts: list[int] = []
        for shape in slide.shapes:
            if shape.left is None:
                continue
            lefts.append(shape.left)
        if len(lefts) < 2:
            continue
        # Cluster left edges; report any pair within 30pt but not within MAX_ALIGN_DRIFT_PT
        sorted_lefts = sorted(set(lefts))
        for i, a in enumerate(sorted_lefts):
            for b in sorted_lefts[i + 1:]:
                diff_pt = emu_to_pt(b - a)
                if diff_pt > 30:
                    break
                if MAX_ALIGN_DRIFT_PT < diff_pt <= 30:
                    violations.append({
                        "slide": slide_idx,
                        "edge": "left",
                        "drift_pt": round(diff_pt, 1),
                        "a_emu": a,
                        "b_emu": b,
                    })
    return violations


def audit_margin(prs, margin_pt: int) -> list[dict]:
    min_emu = margin_pt * PT_TO_EMU
    violations: list[dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not (shape.has_text_frame or getattr(shape, "has_table", False)):
                continue
            if shape.left is None or shape.top is None:
                continue
            if shape.left < min_emu or shape.top < min_emu:
                violations.append({
                    "slide": slide_idx,
                    "shape": shape.name,
                    "left_pt": round(emu_to_pt(shape.left), 1),
                    "top_pt": round(emu_to_pt(shape.top), 1),
                    "min_pt": margin_pt,
                })
    return violations


def audit_empty_slides(prs, plan_layouts: dict[int, str]) -> list[dict]:
    """Section layouts are allowed to have title only; others must have body content."""
    violations: list[dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        layout = plan_layouts.get(slide_idx, "content")
        if layout == "section":
            continue
        text_count = sum(
            1 for _, tf, _bg in collect_text_frames(slide) for p in tf.paragraphs for r in p.runs if r.text.strip()
        )
        if text_count <= 1:  # title only
            violations.append({"slide": slide_idx, "layout": layout, "text_runs": text_count})
    return violations


def main() -> int:
    args = parse_args()
    prs = Presentation(str(args.pptx))
    with args.tokens.open("r", encoding="utf-8") as f:
        tokens = yaml.safe_load(f) or {}

    # Determine theme: prefer tokens.overrides; fall back to template defaults via slide-plan if present
    theme: dict = {}
    template_name = tokens.get("template")
    if template_name:
        try:
            from helpers.templates import load_template
            mod = load_template(template_name)
            theme.update(mod.THEME)
        except Exception:
            pass
    theme.update(tokens.get("overrides", {}))
    margin_pt = int(theme.get("margin_pt", MIN_MARGIN_PT))

    # Load slide layouts from plan if alongside tokens
    plan_layouts: dict[int, str] = {}
    plan_path = args.tokens.parent / "slide-plan.yaml"
    if plan_path.exists():
        plan_data = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
        for i, s in enumerate(plan_data.get("slides", []), start=1):
            plan_layouts[i] = s.get("layout", "content")

    report = {
        "pptx": str(args.pptx),
        "slide_count": len(prs.slides),
        "font_violations": audit_fonts(prs, theme),
        "contrast_violations": audit_contrast(prs, theme, plan_layouts),
        "overflow_violations": audit_overflow(prs),
        "alignment_violations": audit_alignment(prs),
        "margin_violations": audit_margin(prs, margin_pt),
        "empty_slides": audit_empty_slides(prs, plan_layouts),
    }

    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")

    blocker_keys = ("font_violations", "contrast_violations", "overflow_violations", "empty_slides")
    blockers = sum(len(report[k]) for k in blocker_keys)
    warnings = len(report["alignment_violations"]) + len(report["margin_violations"])
    print(f"audit: blockers={blockers} warnings={warnings} → {args.output}")
    return 0 if blockers == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
