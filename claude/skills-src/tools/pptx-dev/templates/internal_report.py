"""Internal report template: navy/gray, dense text-friendly, Yu Gothic UI 18pt."""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.util import Pt

THEME = {
    "name": "internal_report",
    "font_family": "Yu Gothic UI",
    "font_size_title": 28,
    "font_size_body": 18,
    "font_size_notes": 12,
    "color_text": "#1F1F1F",
    "color_title": "#1F4E79",
    "color_accent": "#2E75B6",
    "color_bg": "#FFFFFF",
    "color_section_bg": "#1F4E79",
    "color_section_text": "#FFFFFF",
    "margin_pt": 32,
}


def _apply_theme(run, theme: dict, size: int, color_hex: str) -> None:
    run.font.name = theme["font_family"]
    run.font.size = Pt(size)
    r, g, b = int(color_hex[1:3], 16), int(color_hex[3:5], 16), int(color_hex[5:7], 16)
    run.font.color.rgb = RGBColor(r, g, b)


def render_slide(slide, spec: dict, theme: dict) -> None:
    """Render a single slide spec onto an existing pptx slide. Mutates slide."""
    from pptx.util import Inches

    layout = spec.get("layout", "content")
    title = spec.get("title", "")

    if layout == "title":
        _render_title(slide, spec, theme)
    elif layout == "section":
        _render_section(slide, spec, theme)
    elif layout == "content":
        _render_content(slide, spec, theme)
    elif layout == "two_column":
        _render_two_column(slide, spec, theme)
    elif layout == "chart":
        chart_spec = spec.get("chart", {})
        if chart_spec.get("type") == "diagram":
            _render_diagram(slide, spec, theme)
        else:
            _render_chart_placeholder(slide, spec, theme)
    elif layout == "quote":
        _render_quote(slide, spec, theme)
    elif layout == "closing":
        _render_closing(slide, spec, theme)
    else:
        raise ValueError(f"unknown layout: {layout}")

    notes = spec.get("notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_textbox(slide, left, top, width, height, text, theme, size, color):
    from pptx.util import Emu
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            _apply_theme(run, theme, size, color)
    return box


def _render_title(slide, spec, theme):
    from pptx.util import Inches
    _add_textbox(slide, Inches(0.8).emu, Inches(2.6).emu, Inches(11.8).emu, Inches(1.5).emu,
                 spec.get("title", ""), theme, theme["font_size_title"] + 8, theme["color_title"])
    subtitle = spec.get("subtitle", "")
    if subtitle:
        _add_textbox(slide, Inches(0.8).emu, Inches(4.2).emu, Inches(11.8).emu, Inches(0.8).emu,
                     subtitle, theme, theme["font_size_body"], theme["color_text"])


def _render_section(slide, spec, theme):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    bg = theme["color_section_bg"]
    r, g, b = int(bg[1:3], 16), int(bg[3:5], 16), int(bg[5:7], 16)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
    _add_textbox(slide, Inches(0.8).emu, Inches(3.0).emu, Inches(11.8).emu, Inches(1.5).emu,
                 spec.get("title", ""), theme, theme["font_size_title"] + 4,
                 theme["color_section_text"])


def _render_content(slide, spec, theme):
    from pptx.util import Inches
    _add_textbox(slide, Inches(0.6).emu, Inches(0.5).emu, Inches(12.2).emu, Inches(0.9).emu,
                 spec.get("title", ""), theme, theme["font_size_title"], theme["color_title"])
    bullets = spec.get("bullets", [])
    if not bullets:
        return
    text = "\n".join(f"• {b}" for b in bullets)
    _add_textbox(slide, Inches(0.6).emu, Inches(1.7).emu, Inches(12.2).emu, Inches(5.3).emu,
                 text, theme, theme["font_size_body"], theme["color_text"])


def _render_two_column(slide, spec, theme):
    from pptx.util import Inches
    _add_textbox(slide, Inches(0.6).emu, Inches(0.5).emu, Inches(12.2).emu, Inches(0.9).emu,
                 spec.get("title", ""), theme, theme["font_size_title"], theme["color_title"])
    left = spec.get("left", {})
    right = spec.get("right", {})
    _add_textbox(slide, Inches(0.6).emu, Inches(1.7).emu, Inches(5.9).emu, Inches(0.6).emu,
                 left.get("heading", ""), theme, theme["font_size_body"] + 2, theme["color_accent"])
    _add_textbox(slide, Inches(0.6).emu, Inches(2.4).emu, Inches(5.9).emu, Inches(4.6).emu,
                 "\n".join(f"• {b}" for b in left.get("bullets", [])),
                 theme, theme["font_size_body"], theme["color_text"])
    _add_textbox(slide, Inches(6.8).emu, Inches(1.7).emu, Inches(5.9).emu, Inches(0.6).emu,
                 right.get("heading", ""), theme, theme["font_size_body"] + 2, theme["color_accent"])
    _add_textbox(slide, Inches(6.8).emu, Inches(2.4).emu, Inches(5.9).emu, Inches(4.6).emu,
                 "\n".join(f"• {b}" for b in right.get("bullets", [])),
                 theme, theme["font_size_body"], theme["color_text"])


def _render_chart_placeholder(slide, spec, theme):
    from pptx.util import Inches
    _add_textbox(slide, Inches(0.6).emu, Inches(0.5).emu, Inches(12.2).emu, Inches(0.9).emu,
                 spec.get("title", ""), theme, theme["font_size_title"], theme["color_title"])
    chart = spec.get("chart", {})
    placeholder = f"[chart: {chart.get('type', 'unknown')} source={chart.get('data_source', 'TBD')}]"
    _add_textbox(slide, Inches(0.6).emu, Inches(1.8).emu, Inches(12.2).emu, Inches(4.6).emu,
                 placeholder, theme, theme["font_size_body"], theme["color_text"])
    annotation = chart.get("annotation", "")
    if annotation:
        _add_textbox(slide, Inches(0.6).emu, Inches(6.5).emu, Inches(12.2).emu, Inches(0.5).emu,
                     annotation, theme, theme["font_size_notes"], theme["color_accent"])


def _hex_to_rgbcolor(hex_str):
    r, g, b = int(hex_str[1:3], 16), int(hex_str[3:5], 16), int(hex_str[5:7], 16)
    return RGBColor(r, g, b)


def _add_diagram_box(slide, box, theme):
    """Render one rounded-rectangle box with left-aligned multi-line text.

    box keys: x, y, w, h (inches), lines (list[str]), fill (hex, optional),
    outline (hex, optional), font_size (pt, optional), text_color (hex, optional)
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(box["x"]), Inches(box["y"]), Inches(box["w"]), Inches(box["h"]),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgbcolor(box.get("fill", theme["color_bg"]))
    shape.line.color.rgb = _hex_to_rgbcolor(box.get("outline", theme["color_title"]))
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    font_size = box.get("font_size", theme["font_size_body"])
    text_color = box.get("text_color", theme["color_text"])
    lines = box.get("lines", [])
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        _apply_theme(run, theme, font_size, text_color)
    return shape


def _add_diagram_label(slide, label, theme):
    """Plain caption text box (no fill/border), used as a group heading."""
    from pptx.util import Inches
    _add_textbox(
        slide, Inches(label["x"]).emu, Inches(label["y"]).emu,
        Inches(label["w"]).emu, Inches(label.get("h", 0.3)).emu,
        label["text"], theme, label.get("font_size", theme["font_size_body"] - 1),
        label.get("color", theme["color_accent"]),
    )


def _add_diagram_arrow(slide, from_box, to_box, theme):
    """Straight connector with a triangle arrowhead from from_box's right-center
    to to_box's left-center."""
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    start_x = from_box["x"] + from_box["w"]
    start_y = from_box["y"] + from_box["h"] / 2
    end_x = to_box["x"]
    end_y = to_box["y"] + to_box["h"] / 2

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(start_x), Inches(start_y), Inches(end_x), Inches(end_y)
    )
    connector.line.color.rgb = _hex_to_rgbcolor(theme["color_title"])
    connector.line.width = Pt(1.5)
    ln = connector.line._get_or_add_ln()
    tail_end = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail_end)
    return connector


def _set_cell_text(cell, text, theme, font_size, color):
    """Set a table cell's text via an explicit run (cell.text = "" leaves the
    paragraph with zero runs, so a later runs[0] lookup raises IndexError)."""
    tf = cell.text_frame
    tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = text
    _apply_theme(run, theme, font_size, color)


def _add_diagram_table(slide, table, theme):
    """Render a native pptx table (grid) for concrete before/after examples.

    table keys: id, x, y, w, h (inches), headers: [str], rows: [[str]],
    font_size (optional), highlight_rows: [int] (0-based row index into `rows`,
    not counting the header), highlight_fill (hex, optional)
    """
    from pptx.util import Inches, Pt

    headers = table.get("headers", [])
    rows = table.get("rows", [])
    n_cols = len(headers)
    n_rows = len(rows) + 1  # + header row

    if n_cols == 0:
        # pptx internally divides table width by column count, so 0 columns
        # crashes with a bare ZeroDivisionError — fail with a clearer cause.
        raise ValueError(f"diagram table {table.get('id', '?')!r} has no headers (column count must be >= 1)")

    graphic_frame = slide.shapes.add_table(
        n_rows, n_cols, Inches(table["x"]), Inches(table["y"]), Inches(table["w"]), Inches(table["h"])
    )
    pptx_table = graphic_frame.table
    font_size = table.get("font_size", theme["font_size_body"] - 3)
    highlight_rows = set(table.get("highlight_rows", []))
    highlight_fill = table.get("highlight_fill", "#FFF3CD")

    for c, header_text in enumerate(headers):
        cell = pptx_table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgbcolor(theme["color_title"])
        _set_cell_text(cell, header_text, theme, font_size, "#FFFFFF")

    for r, row_values in enumerate(rows):
        for c, value in enumerate(row_values):
            cell = pptx_table.cell(r + 1, c)
            if r in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex_to_rgbcolor(highlight_fill)
            _set_cell_text(cell, str(value), theme, font_size, theme["color_text"])

    return graphic_frame


def _render_diagram(slide, spec, theme):
    """Render a box-and-arrow architecture diagram.

    spec.chart.diagram keys:
      boxes: [{id, x, y, w, h, lines, fill, outline, font_size, text_color}]
      tables: [{id, x, y, w, h, headers, rows, font_size, highlight_rows, highlight_fill}]
      arrows: [{from: box_or_table_id, to: box_or_table_id}]
      labels: [{text, x, y, w, h, font_size, color}]   # plain captions, no box
      footnote: [str, ...]                              # bottom caption lines
    """
    from pptx.util import Inches

    _add_textbox(slide, Inches(0.6).emu, Inches(0.45).emu, Inches(12.2).emu, Inches(0.65).emu,
                 spec.get("title", ""), theme, theme["font_size_title"], theme["color_title"])

    subtitle = spec.get("subtitle", "")
    if subtitle:
        _add_textbox(slide, Inches(0.6).emu, Inches(1.15).emu, Inches(12.2).emu, Inches(0.35).emu,
                     subtitle, theme, theme["font_size_body"] - 2, theme["color_text"])

    diagram = spec.get("chart", {}).get("diagram", {})
    boxes_by_id = {}
    for box in diagram.get("boxes", []):
        _add_diagram_box(slide, box, theme)
        boxes_by_id[box["id"]] = box

    for table in diagram.get("tables", []):
        _add_diagram_table(slide, table, theme)
        boxes_by_id[table["id"]] = table

    for label in diagram.get("labels", []):
        _add_diagram_label(slide, label, theme)

    for arrow in diagram.get("arrows", []):
        from_box = boxes_by_id[arrow["from"]]
        to_box = boxes_by_id[arrow["to"]]
        _add_diagram_arrow(slide, from_box, to_box, theme)

    footnote_lines = diagram.get("footnote", [])
    for i, line in enumerate(footnote_lines):
        _add_textbox(slide, Inches(0.6).emu, Inches(6.35 + i * 0.35).emu,
                     Inches(12.2).emu, Inches(0.35).emu,
                     line, theme, theme["font_size_body"] - 2, theme["color_accent"])


def _render_quote(slide, spec, theme):
    from pptx.util import Inches
    quote = spec.get("quote", "")
    source = spec.get("source", "")
    _add_textbox(slide, Inches(1.5).emu, Inches(2.5).emu, Inches(10.3).emu, Inches(2.5).emu,
                 f"“{quote}”", theme, theme["font_size_title"], theme["color_text"])
    if source:
        _add_textbox(slide, Inches(1.5).emu, Inches(5.2).emu, Inches(10.3).emu, Inches(0.6).emu,
                     f"— {source}", theme, theme["font_size_body"], theme["color_accent"])


def _render_closing(slide, spec, theme):
    from pptx.util import Inches
    _add_textbox(slide, Inches(0.8).emu, Inches(2.8).emu, Inches(11.8).emu, Inches(1.4).emu,
                 spec.get("title", "Thank you"), theme, theme["font_size_title"] + 12,
                 theme["color_title"])
    cta = spec.get("cta", "")
    if cta:
        _add_textbox(slide, Inches(0.8).emu, Inches(4.3).emu, Inches(11.8).emu, Inches(1.0).emu,
                     cta, theme, theme["font_size_body"] + 2, theme["color_accent"])
