#!/usr/bin/env python3
"""Generate a .pptx from slide-plan.yaml + tokens.yaml + a template."""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

import yaml
from pptx import Presentation

# Allow importing from sibling templates/ directory
sys.path.insert(0, str(Path(__file__).resolve().parent))

from helpers.templates import load_template  # noqa: E402


VALID_LAYOUTS = {"title", "section", "content", "two_column", "chart", "quote", "closing"}


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Generate .pptx from slide plan + tokens")
    p.add_argument("--plan", type=Path, required=True, help="path to slide-plan.yaml")
    p.add_argument("--tokens", type=Path, required=True, help="path to tokens.yaml")
    p.add_argument("--output", type=Path, required=True, help="output .pptx path")
    return p.parse_args()


def load_yaml(path: Path) -> dict:
    if not path.exists():
        raise FileNotFoundError(f"file not found: {path}")
    with path.open("r", encoding="utf-8") as f:
        data = yaml.safe_load(f)
    if not isinstance(data, dict):
        raise ValueError(f"{path}: expected mapping at root")
    return data


def validate_plan(plan: dict) -> list[dict]:
    if "slides" not in plan or not isinstance(plan["slides"], list):
        raise ValueError("slide-plan.yaml must have a `slides` list")
    slides = plan["slides"]
    if not slides:
        raise ValueError("slide-plan.yaml `slides` is empty")
    seen_ids: set[int] = set()
    for i, s in enumerate(slides):
        if not isinstance(s, dict):
            raise ValueError(f"slide #{i}: expected mapping")
        sid = s.get("id")
        if not isinstance(sid, int):
            raise ValueError(f"slide #{i}: `id` must be int")
        if sid in seen_ids:
            raise ValueError(f"duplicate slide id: {sid}")
        seen_ids.add(sid)
        layout = s.get("layout")
        if layout not in VALID_LAYOUTS:
            raise ValueError(f"slide id={sid}: invalid layout `{layout}`; expected one of {VALID_LAYOUTS}")
    return slides


def merge_theme(template_theme: dict, tokens: dict) -> dict:
    """Apply tokens.overrides onto template defaults."""
    merged = dict(template_theme)
    overrides = tokens.get("overrides", {})
    if overrides and not isinstance(overrides, dict):
        raise ValueError("tokens.yaml `overrides` must be a mapping")
    if overrides:
        # Translate user-facing keys to theme keys
        mapping = {
            "accent_color": "color_accent",
            "text_color": "color_text",
            "title_color": "color_title",
            "bg_color": "color_bg",
            "font_family": "font_family",
            "font_size_body": "font_size_body",
            "font_size_title": "font_size_title",
            "margin_pt": "margin_pt",
        }
        for user_key, value in overrides.items():
            theme_key = mapping.get(user_key, user_key)
            merged[theme_key] = value
    return merged


def main() -> int:
    args = parse_args()
    plan = load_yaml(args.plan)
    tokens = load_yaml(args.tokens)

    template_name = tokens.get("template") or plan.get("template")
    if not template_name:
        raise ValueError("template name not found in tokens.yaml or slide-plan.yaml")

    template = load_template(template_name)
    theme = merge_theme(template.THEME, tokens)
    slides_spec = validate_plan(plan)

    prs = Presentation()
    # 16:9
    prs.slide_width = 12192000
    prs.slide_height = 6858000
    blank_layout = prs.slide_layouts[6]  # blank

    for spec in slides_spec:
        slide = prs.slides.add_slide(blank_layout)
        template.render_slide(slide, spec, theme)

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(f"generated: {args.output} ({len(slides_spec)} slides, template={template_name})")
    return 0


if __name__ == "__main__":
    sys.exit(main())
